from __future__ import annotations

import argparse
import hashlib
import json
import os
import re
import sys
import time
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Callable

from dotenv import load_dotenv
import PIL.Image
import io

try:
    from google import genai
    from google.genai import types
    HAS_GENAI = True
except ImportError:
    HAS_GENAI = False

try:
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    HAS_PPTX_ENUMS = True
except ImportError:
    HAS_PPTX_ENUMS = False

try:
    from openai import OpenAI
    HAS_OPENAI = True
except ImportError:
    HAS_OPENAI = False


# ── Azure AI Foundry defaults (overridden at runtime via env vars) ─────────────
_AZURE_ENDPOINT_DEFAULT = "https://haulv226038-3382-resource.services.ai.azure.com/openai/v1"
_AZURE_DEPLOYMENT_DEFAULT = "gpt-4-1-mini-2025-04-14-ft-cd68c6dbd12543298c5eb1ad64af5e4a"

# ── Groq defaults (tertiary fallback) ─────────────────────────────────────────
_GROQ_BASE_URL = "https://api.groq.com/openai/v1"
_GROQ_MODEL_DEFAULT = "meta-llama/llama-4-scout-17b-16e-instruct"

# Module-level LLM call log. Reset by run_mode at the start of each experiment
# run so inventory/gap-report calls don't pollute experiment records.
_llm_call_records: list[dict[str, Any]] = []
_progress_callback: Callable[[dict[str, Any]], None] | None = None
_progress_records: list[dict[str, Any]] = []


def _record_llm_call(
    operation: str,
    provider: str,
    model: str,
    status: str,
    latency_ms: int,
    prompt_tokens: int = 0,
    completion_tokens: int = 0,
    error: str | None = None,
) -> None:
    input_rate = float(os.getenv("AZURE_OPENAI_INPUT_COST_PER_1M", "0.15"))
    output_rate = float(os.getenv("AZURE_OPENAI_OUTPUT_COST_PER_1M", "0.60"))
    # Cost applies only to Azure calls; Gemini records cost=0 (separate billing).
    if provider == "azure_openai":
        cost = (prompt_tokens / 1_000_000) * input_rate + (completion_tokens / 1_000_000) * output_rate
    else:
        cost = 0.0
    record: dict[str, Any] = {
        "operation": operation,
        "provider": provider,
        "model": model,
        "status": status,
        "latency_ms": latency_ms,
        "prompt_tokens": prompt_tokens,
        "completion_tokens": completion_tokens,
        "total_tokens": prompt_tokens + completion_tokens,
        "estimated_cost_usd": round(cost, 8),
    }
    if error:
        record["error"] = error
    _llm_call_records.append(record)


def _summarize_llm_metrics() -> dict[str, Any]:
    records = list(_llm_call_records)
    total = len(records)
    successful = sum(1 for r in records if r["status"] == "success")
    latencies = [r["latency_ms"] for r in records]
    return {
        "total_calls": total,
        "successful_calls": successful,
        "failed_calls": total - successful,
        "total_prompt_tokens": sum(r["prompt_tokens"] for r in records),
        "total_completion_tokens": sum(r["completion_tokens"] for r in records),
        "total_tokens": sum(r["total_tokens"] for r in records),
        "total_cost_usd": round(sum(r["estimated_cost_usd"] for r in records), 8),
        "avg_latency_ms": round(sum(latencies) / max(len(latencies), 1), 1),
        "records": records,
    }


SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".jpg", ".jpeg", ".png"}
IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png"}

EXPECTED_TYPES = {"text_pdf", "scanned_pdf", "slide_pdf", "mixed_pdf", "pptx", "image"}
EXPECTED_ROUTES = {"direct_text", "hybrid", "vision_only"}

REQUIRED_EDGE_CASES = [
    "text_pdf",
    "scanned_pdf",
    "slide_pdf",
    "mixed_pdf",
    "pptx",
    "image",
]


def utc_now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()


def _emit_progress(event: str, **fields: Any) -> None:
    payload = {"timestamp": utc_now_iso(), "event": event, **fields}
    _progress_records.append(payload)
    if _progress_callback is None:
        return
    try:
        _progress_callback(payload)
    except Exception:
        pass


def file_sha256(path: Path) -> str:
    sha = hashlib.sha256()
    with path.open("rb") as handle:
        for block in iter(lambda: handle.read(1024 * 1024), b""):
            sha.update(block)
    return sha.hexdigest()


def ensure_dir(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


def read_json(path: Path) -> Any:
    with path.open("r", encoding="utf-8") as handle:
        return json.load(handle)


def write_json(path: Path, payload: Any) -> None:
    ensure_dir(path.parent)
    with path.open("w", encoding="utf-8") as handle:
        json.dump(payload, handle, ensure_ascii=False, indent=2)
        handle.write("\n")


def write_text(path: Path, content: str) -> None:
    ensure_dir(path.parent)
    with path.open("w", encoding="utf-8") as handle:
        handle.write(content)


def resolve_project_root(explicit_root: str | None) -> Path:
    if explicit_root:
        return Path(explicit_root).resolve()
    return Path(__file__).resolve().parents[3]


def resolve_output_root(project_root: Path) -> Path:
    configured = os.getenv("DOCUMENT_PROCESSING_OUTPUT_BASE", "data/pipeline_outputs/document_processing").strip()
    base = Path(configured)
    if not base.is_absolute():
        base = project_root / base
    return base


def discover_samples(sample_root: Path) -> list[Path]:
    if not sample_root.exists():
        return []

    discovered: list[Path] = []
    for path in sample_root.rglob("*"):
        if path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS:
            discovered.append(path)
    return sorted(discovered)


def parse_sample_metadata(project_root: Path, file_path: Path) -> dict[str, Any]:
    rel = file_path.relative_to(project_root).as_posix()
    parts = rel.split("/")

    tier = "unknown"
    course_code = "unknown"
    doc_subtype = "unknown"

    if len(parts) >= 6 and parts[0] == "data" and parts[1] == "sample":
        tier = parts[2]
        course_code = parts[3]
        doc_subtype = parts[4]

    ext = file_path.suffix.lower()
    if ext == ".pdf":
        format_name = "pdf"
    elif ext == ".pptx":
        format_name = "pptx"
    elif ext in IMAGE_EXTENSIONS:
        format_name = "image"
    else:
        format_name = "other"

    return {
        "path": rel,
        "tier": tier,
        "course_code": course_code,
        "doc_subtype": doc_subtype,
        "extension": ext,
        "format": format_name,
        "file_name": file_path.name,
        "size_bytes": file_path.stat().st_size,
    }


def classify_pdf_type(pdf_path: Path) -> tuple[str, dict[str, Any]]:
    try:
        import fitz  # PyMuPDF
    except Exception:
        return "text_pdf", {
            "reason": "pymupdf_unavailable",
            "text_chars": None,
            "page_count": None,
            "image_count": None,
        }

    doc = fitz.open(pdf_path)
    page_count = len(doc)
    text_chars = 0
    image_count = 0
    landscape_count = 0

    for page in doc:
        text_chars += len((page.get_text("text") or "").strip())
        image_count += len(page.get_images(full=True))
        if page.rect.width > page.rect.height:
            landscape_count += 1

    doc.close()
    avg_chars_per_page = (text_chars / page_count) if page_count else 0
    images_per_page = (image_count / page_count) if page_count else 0
    is_landscape = (landscape_count / page_count) > 0.5 if page_count else False

    if text_chars < 80:
        guessed_type = "scanned_pdf"
    elif is_landscape and image_count > 0:
        guessed_type = "slide_pdf"
    elif avg_chars_per_page >= 120 and images_per_page <= 0.3:
        guessed_type = "text_pdf"
    elif image_count > 0 and avg_chars_per_page >= 80:
        guessed_type = "mixed_pdf"
    else:
        guessed_type = "text_pdf"

    evidence = {
        "text_chars": text_chars,
        "page_count": page_count,
        "image_count": image_count,
        "avg_chars_per_page": round(avg_chars_per_page, 2),
        "images_per_page": round(images_per_page, 3),
        "landscape_pages": landscape_count,
        "is_landscape": is_landscape,
    }
    return guessed_type, evidence


def inventory_mode(project_root: Path, output_root: Path) -> dict[str, Any]:
    sample_root = project_root / "data" / "sample"
    files = discover_samples(sample_root)

    by_course_tier: dict[str, dict[str, dict[str, int]]] = {}
    format_counts = {"pdf": 0, "pptx": 0, "image": 0, "other": 0}
    inferred_types: dict[str, int] = {
        "text_pdf": 0,
        "scanned_pdf": 0,
        "mixed_pdf": 0,
        "pptx": 0,
        "image": 0,
    }
    sample_items: list[dict[str, Any]] = []

    for file_path in files:
        item = parse_sample_metadata(project_root, file_path)
        format_counts[item["format"]] = format_counts.get(item["format"], 0) + 1

        key = f"{item['tier']}:{item['course_code']}"
        by_course_tier.setdefault(key, {}).setdefault(item["doc_subtype"], {})
        by_course_tier[key][item["doc_subtype"]][item["format"]] = (
            by_course_tier[key][item["doc_subtype"]].get(item["format"], 0) + 1
        )

        if item["format"] == "pdf":
            inferred_type, evidence = classify_pdf_type(file_path)
            inferred_types[inferred_type] += 1
            item["inferred_expected_type"] = inferred_type
            item["type_evidence"] = evidence
        elif item["format"] == "pptx":
            inferred_types["pptx"] += 1
            item["inferred_expected_type"] = "pptx"
        elif item["format"] == "image":
            inferred_types["image"] += 1
            item["inferred_expected_type"] = "image"

        sample_items.append(item)

    report = {
        "generated_at": utc_now_iso(),
        "sample_root": sample_root.relative_to(project_root).as_posix(),
        "total_files": len(files),
        "format_counts": format_counts,
        "inferred_edge_case_counts": inferred_types,
        "by_course_tier": by_course_tier,
        "files": sample_items,
    }

    write_json(output_root / "inventory_report.json", report)
    return report


def gap_report_mode(project_root: Path, output_root: Path) -> dict[str, Any]:
    inventory_path = output_root / "inventory_report.json"
    if inventory_path.exists():
        inventory = read_json(inventory_path)
    else:
        inventory = inventory_mode(project_root, output_root)

    inferred_counts = inventory.get("inferred_edge_case_counts", {})
    coverage: list[dict[str, Any]] = []
    missing: list[str] = []

    for category in REQUIRED_EDGE_CASES:
        count = int(inferred_counts.get(category, 0))
        status = "covered" if count > 0 else "missing"
        if status == "missing":
            missing.append(category)
        coverage.append({
            "category": category,
            "status": status,
            "sample_count": count,
        })

    report = {
        "generated_at": utc_now_iso(),
        "required_categories": REQUIRED_EDGE_CASES,
        "coverage": coverage,
        "missing_categories": missing,
        "policy": {
            "auto_synthesis_allowed": False,
            "instruction": "Add minimum manual samples only for missing categories.",
        },
    }

    write_json(output_root / "gap_report.json", report)
    return report


def expected_type_to_route(expected_type: str) -> str:
    mapping = {
        "text_pdf": "direct_text",
        "scanned_pdf": "vision_only",
        "slide_pdf": "vision_only",
        "mixed_pdf": "hybrid",
        "pptx": "hybrid",
        "image": "vision_only",
    }
    return mapping.get(expected_type, "hybrid")


def decide_route(input_file: Path) -> tuple[str, dict[str, Any]]:
    ext = input_file.suffix.lower()
    if ext in IMAGE_EXTENSIONS:
        return "vision_only", {"reason": "image_input"}
    if ext == ".pptx":
        return "hybrid", {"reason": "pptx_text_plus_embedded_images"}

    if ext == ".pdf":
        inferred_type, evidence = classify_pdf_type(input_file)
        return expected_type_to_route(inferred_type), {
            "reason": "pdf_probe",
            "inferred_type": inferred_type,
            **evidence,
        }

    return "hybrid", {"reason": "fallback"}


def _call_gemini_vlm(
    prompt: str,
    image_data: bytes | list[bytes] | None = None,
    operation: str = "vision",
) -> str:
    """Call Gemini for vision. Raises on failure so caller can fall back."""
    if not HAS_GENAI:
        raise RuntimeError("google-genai not installed")
    api_key = os.getenv("GEMINI_API_KEY")
    if not api_key:
        raise RuntimeError("GEMINI_API_KEY_missing")

    model_id = os.getenv("VLM_MODEL_ID", "gemini-3.1-flash-lite")
    client = genai.Client(api_key=api_key, http_options={"timeout": 30})
    t0 = time.time()
    _emit_progress(
        "llm_call_start",
        operation=operation,
        provider="gemini",
        model=model_id,
        has_image=bool(image_data),
    )
    try:
        contents: list[Any] = []
        if image_data:
            imgs = image_data if isinstance(image_data, list) else [image_data]
            for img in imgs:
                contents.append(PIL.Image.open(io.BytesIO(img)))
        contents.append(prompt)
        response = client.models.generate_content(
            model=model_id,
            contents=contents,
            config=types.GenerateContentConfig(max_output_tokens=4096, temperature=0.3),
        )
        latency_ms = int((time.time() - t0) * 1000)
        usage = response.usage_metadata
        _record_llm_call(
            operation, "gemini", model_id, "success", latency_ms,
            prompt_tokens=getattr(usage, "prompt_token_count", 0) or 0,
            completion_tokens=getattr(usage, "candidates_token_count", 0) or 0,
        )
        _emit_progress(
            "llm_call_success",
            operation=operation,
            provider="gemini",
            model=model_id,
            latency_ms=latency_ms,
        )
        return response.text or ""
    except Exception as e:
        latency_ms = int((time.time() - t0) * 1000)
        status_code = getattr(e, "status_code", None) or getattr(e, "code", None)
        err_detail = f"{type(e).__name__}(HTTP_{status_code}): {e}" if status_code else f"{type(e).__name__}: {e}"
        _record_llm_call(operation, "gemini", model_id, "error", latency_ms, error=err_detail)
        _emit_progress(
            "llm_call_error",
            operation=operation,
            provider="gemini",
            model=model_id,
            latency_ms=latency_ms,
            error=err_detail,
        )
        raise


def _call_azure_openai(
    prompt: str,
    image_data: bytes | list[bytes] | None = None,
    original_exc: Exception | None = None,
    operation: str = "text",
) -> str:
    """Call Azure AI Foundry for vision fallback or text normalization. Records metrics."""
    if not HAS_OPENAI:
        msg = "openai_client_not_installed"
        _record_llm_call(operation, "azure_openai", _AZURE_DEPLOYMENT_DEFAULT, "error", 0, error=msg)
        if original_exc:
            return f"[VISION_ERROR] Gemini failed ({original_exc}) and Azure fallback unavailable ({msg})."
        return f"[LLM_ERROR] {msg}."

    endpoint = os.getenv("AZURE_OPENAI_ENDPOINT", _AZURE_ENDPOINT_DEFAULT)
    deployment = os.getenv("AZURE_OPENAI_DEPLOYMENT", _AZURE_DEPLOYMENT_DEFAULT)
    api_key = os.getenv("AZURE_AI_API_KEY")
    if not api_key:
        msg = "AZURE_AI_API_KEY_missing"
        _record_llm_call(operation, "azure_openai", deployment, "error", 0, error=msg)
        if original_exc:
            return f"[VISION_ERROR] Gemini failed ({original_exc}) and Azure unavailable ({msg})."
        return f"[LLM_ERROR] {msg}."

    for attempt in range(2):
        t0 = time.time()
        _emit_progress(
            "llm_call_start",
            operation=operation,
            provider="azure_openai",
            model=deployment,
            attempt=attempt + 1,
            has_image=bool(image_data),
        )
        try:
            client = OpenAI(
                base_url=endpoint,
                api_key=api_key,
                timeout=30.0,
            )

            if image_data:
                import base64
                content_parts: list[Any] = [{"type": "text", "text": prompt}]
                imgs = image_data if isinstance(image_data, list) else [image_data]
                for img in imgs:
                    b64 = base64.b64encode(img).decode("utf-8")
                    content_parts.append({
                        "type": "image_url",
                        "image_url": {"url": f"data:image/jpeg;base64,{b64}"},
                    })
                messages: list[Any] = [{"role": "user", "content": content_parts}]
            else:
                messages = [{"role": "user", "content": prompt}]

            response = client.chat.completions.create(
                model=deployment,
                messages=messages,
                temperature=0.3,
                max_tokens=4096,
            )

            latency_ms = int((time.time() - t0) * 1000)
            usage = response.usage
            _record_llm_call(
                operation, "azure_openai", deployment, "success", latency_ms,
                prompt_tokens=usage.prompt_tokens if usage else 0,
                completion_tokens=usage.completion_tokens if usage else 0,
            )
            _emit_progress(
                "llm_call_success",
                operation=operation,
                provider="azure_openai",
                model=deployment,
                attempt=attempt + 1,
                latency_ms=latency_ms,
            )
            return response.choices[0].message.content or ""

        except Exception as e:
            latency_ms = int((time.time() - t0) * 1000)
            status_code = getattr(e, "status_code", None)
            err_detail = f"{type(e).__name__}(HTTP_{status_code}): {e}" if status_code else f"{type(e).__name__}: {e}"
            _record_llm_call(operation, "azure_openai", deployment, "error", latency_ms, error=err_detail)
            _emit_progress(
                "llm_call_error",
                operation=operation,
                provider="azure_openai",
                model=deployment,
                attempt=attempt + 1,
                latency_ms=latency_ms,
                error=err_detail,
            )
            if attempt == 0:
                time.sleep(2.0)
            else:
                if original_exc:
                    return f"[VISION_ERROR] Gemini failed ({original_exc}) and Azure fallback failed twice ({err_detail})."
                return f"[LLM_ERROR] Azure OpenAI failed twice ({err_detail})."
    return "[LLM_ERROR] unreachable"


def _call_groq(
    prompt: str,
    image_data: bytes | list[bytes] | None = None,
    original_exc: Exception | None = None,
    operation: str = "text",
) -> str:
    """Call Groq as final vision/text fallback. Caller must set GROQ_MODEL_ID to a vision-capable model when image_data is passed."""
    if not HAS_OPENAI:
        msg = "openai_client_not_installed"
        _record_llm_call(operation, "groq", _GROQ_MODEL_DEFAULT, "error", 0, error=msg)
        return f"[LLM_ERROR] Groq unavailable ({msg})."
    api_key = os.getenv("GROQ_API_KEY")
    if not api_key:
        msg = "GROQ_API_KEY_missing"
        _record_llm_call(operation, "groq", _GROQ_MODEL_DEFAULT, "error", 0, error=msg)
        return f"[LLM_ERROR] Groq unavailable ({msg})."

    model = os.getenv("GROQ_MODEL_ID", _GROQ_MODEL_DEFAULT)
    t0 = time.time()
    _emit_progress(
        "llm_call_start",
        operation=operation,
        provider="groq",
        model=model,
        has_image=bool(image_data),
    )
    try:
        client = OpenAI(api_key=api_key, base_url=_GROQ_BASE_URL, timeout=30.0)
        if image_data:
            import base64
            content_parts: list[Any] = [{"type": "text", "text": prompt}]
            imgs = image_data if isinstance(image_data, list) else [image_data]
            for img in imgs:
                b64 = base64.b64encode(img).decode("utf-8")
                content_parts.append({
                    "type": "image_url",
                    "image_url": {"url": f"data:image/jpeg;base64,{b64}"},
                })
            messages: list[Any] = [{"role": "user", "content": content_parts}]
        else:
            messages = [{"role": "user", "content": prompt}]
        response = client.chat.completions.create(
            model=model,
            messages=messages,
            temperature=0.3,
            max_tokens=4096,
        )
        latency_ms = int((time.time() - t0) * 1000)
        usage = response.usage
        _record_llm_call(
            operation, "groq", model, "success", latency_ms,
            prompt_tokens=usage.prompt_tokens if usage else 0,
            completion_tokens=usage.completion_tokens if usage else 0,
        )
        _emit_progress(
            "llm_call_success",
            operation=operation,
            provider="groq",
            model=model,
            latency_ms=latency_ms,
        )
        return response.choices[0].message.content or ""
    except Exception as e:
        latency_ms = int((time.time() - t0) * 1000)
        status_code = getattr(e, "status_code", None)
        err_detail = f"{type(e).__name__}(HTTP_{status_code}): {e}" if status_code else f"{type(e).__name__}: {e}"
        _record_llm_call(operation, "groq", model, "error", latency_ms, error=err_detail)
        _emit_progress(
            "llm_call_error",
            operation=operation,
            provider="groq",
            model=model,
            latency_ms=latency_ms,
            error=err_detail,
        )
        if original_exc:
            return f"[LLM_ERROR] All providers failed: Gemini ({original_exc}), Azure, Groq ({err_detail})."
        return f"[LLM_ERROR] Azure and Groq both failed: Groq error ({err_detail})."


def call_vlm(prompt: str, image_data: bytes | list[bytes] | None = None) -> str:
    """Vision call: Azure (primary) → Gemini (fallback) → Groq (last resort, vision-capable model required)."""
    if not os.getenv("OCR_ENABLE_REAL_VISION", "false").lower() == "true":
        return "[VISION_PLACEHOLDER] Real vision disabled."

    azure_result = _call_azure_openai(prompt, image_data, operation="vision")
    if not azure_result.startswith(("[LLM_ERROR]", "[VISION_ERROR]")):
        return azure_result

    gemini_exc: Exception | None = None
    if HAS_GENAI and os.getenv("GEMINI_API_KEY"):
        try:
            return _call_gemini_vlm(prompt, image_data, operation="vision_fallback_gemini")
        except Exception as e:
            gemini_exc = e

    return _call_groq(prompt, image_data=image_data, original_exc=gemini_exc, operation="vision_fallback_groq")


# ── Visual classification ─────────────────────────────────────────────────────

VISUAL_CATEGORIES = ["table_or_matrix", "graph_diagram", "formula", "general_visual", "decorative"]

_LEARNING_VALUE_BY_LABEL = {
    "table_or_matrix": "high",
    "graph_diagram": "high",
    "formula": "high",
    "general_visual": "medium",
    "decorative": "none",
}


def _make_classification(
    label: str,
    learning_value: str,
    confidence: float,
    reason: str,
    position_hint: str = "inline",
) -> dict[str, Any]:
    """Build a classification dict and determine the action from label + context.

    Action rules:
    - decorative + (header/footer position OR confidence >= 0.75) → skip
    - decorative + inline + confidence < 0.75 → minimal_tag
      (uncertain decorative inline might be a meaningful figure; tag rather than drop)
    - all other labels → extract
    """
    if label == "decorative":
        if confidence >= 0.75 or position_hint in ("header", "footer"):
            action = "skip"
        else:
            action = "minimal_tag"
    else:
        action = "extract"

    return {
        "label": label,
        "learning_value": learning_value,
        "action": action,
        "confidence": round(float(confidence), 3),
        "reason": reason,
    }


def classify_visual(
    image_bytes: bytes,
    surrounding_text: str = "",
    position_hint: str = "inline",
) -> dict[str, Any]:
    """Classify an image using the VLM with surrounding context.

    Returns a dict: {label, learning_value, action, confidence, reason}
    - label: one of VISUAL_CATEGORIES
    - action: "skip" | "minimal_tag" | "extract"
    - confidence: 0.0–1.0

    Falls back gracefully when vision is disabled or VLM response is not valid JSON.
    """
    if not os.getenv("OCR_ENABLE_REAL_VISION", "false").lower() == "true":
        return _make_classification("general_visual", "medium", 0.5, "vision_disabled", position_hint)

    context_snippet = (surrounding_text or "").strip()[:1500]
    prompt = (
        "Classify this image and return a JSON object with exactly these four fields:\n"
        '  "label": one of [table_or_matrix, graph_diagram, formula, general_visual, decorative]\n'
        '  "learning_value": one of [none, low, medium, high]\n'
        '  "confidence": float 0.0-1.0\n'
        '  "reason": one short sentence\n\n'
        "Category guide:\n"
        "- table_or_matrix: data grids, matrices, structured tables\n"
        "- graph_diagram: nodes/edges, flowcharts, system architecture, conceptual diagrams (high learning value)\n"
        "- formula: equations, symbolic math, LaTeX-style notation\n"
        "- general_visual: figures with academic content — charts, screenshots that teach something\n"
        "- decorative: logos, course/school branding, page borders, repeated banners, dividers,\n"
        "  background flourishes, page-number ornaments, watermarks. NO learning value.\n\n"
        f"Position on page: {position_hint}\n"
        f"Surrounding text:\n---\n{context_snippet}\n---\n\n"
        'Return ONLY valid JSON, e.g. {"label":"decorative","learning_value":"none","confidence":0.9,"reason":"logo with no academic content"}'
    )
    raw = call_vlm(prompt, image_bytes).strip()

    # Try JSON parse first.
    try:
        match = re.search(r"\{[^{}]+\}", raw, re.DOTALL)
        if match:
            parsed = json.loads(match.group(0))
            label = parsed.get("label", "general_visual")
            if label not in VISUAL_CATEGORIES:
                label = "general_visual"
            learning_value = parsed.get("learning_value", _LEARNING_VALUE_BY_LABEL.get(label, "medium"))
            confidence = float(parsed.get("confidence", 0.5))
            reason = str(parsed.get("reason", "vlm_classified"))
            return _make_classification(label, learning_value, confidence, reason, position_hint)
    except Exception:
        pass

    # Fallback: scan raw text for a category name.
    lower = raw.lower()
    for cat in VISUAL_CATEGORIES:
        if cat in lower:
            lv = _LEARNING_VALUE_BY_LABEL.get(cat, "medium")
            return _make_classification(cat, lv, 0.5, "fallback_text_parse", position_hint)

    return _make_classification("general_visual", "medium", 0.5, "fallback_default", position_hint)


def _image_position_hint(rect_y0: float | None, rect_y1: float | None, page_height: float | None) -> str:
    if rect_y0 is None or rect_y1 is None or not page_height:
        return "inline"
    top_band = page_height * 0.15
    bottom_band = page_height * 0.85
    if rect_y1 <= top_band:
        return "header"
    if rect_y0 >= bottom_band:
        return "footer"
    return "inline"


def get_specialized_prompt(category: str) -> str:
    if category == "table_or_matrix":
        return "Extract row labels, column labels, and cell values. Return structured JSON with schema: {schema_version, content_type, row_labels, column_labels, values, notes}. Preserve exact numeric values."
    elif category == "graph_diagram":
        return "Extract node labels, edge list and directionality. Return a structured graph representation in JSON."
    elif category == "formula":
        return "Extract surrounding explanatory text separately from formula text. Return formulas in LaTeX."
    elif category == "decorative":
        return "This is a decorative image. Describe it minimally."
    else:
        return "Describe this visual element, focusing only on content that contributes to learning value. If it is a diagram, flowchart, or memory layout, draw/represent it visually using a Markdown table or a text-based ASCII diagram (using boxes [+---+] and arrows [-->] inside a preformatted ``` block) first. Avoid filler commentary about style."


# ── Rule-based cleanup ────────────────────────────────────────────────────────

def rule_based_cleanup(blocks: list[dict[str, Any]]) -> list[dict[str, Any]]:
    cleaned = []
    for b in blocks:
        if b["kind"] == "text":
            text = b.get("content", "")
            text = text.replace("Here's a description", "").replace("Here is the table", "").replace("```json", "").replace("```", "")
            if "image_url" in text:
                text = text.replace("image_url", "")
            b_new = dict(b)
            b_new["content"] = text
            cleaned.append(b_new)
        else:
            cleaned.append(b)
    return cleaned


# ── Normalization ─────────────────────────────────────────────────────────────

NORMALIZATION_CHAR_BUDGET = int(os.getenv("DOCUMENT_PROCESSING_NORMALIZE_CHAR_BUDGET", "24000"))


def batch_blocks_for_normalization(blocks: list[dict[str, Any]], char_budget: int = NORMALIZATION_CHAR_BUDGET) -> list[list[dict[str, Any]]]:
    """Group sequential normalizable blocks into batches that fit a per-call char budget.

    Boundary preference order: slide > page > char budget. A new batch starts when adding the next
    block would exceed the budget, or when crossing a slide boundary (PPTX is naturally self-contained).
    """
    batches: list[list[dict[str, Any]]] = []
    current: list[dict[str, Any]] = []
    current_len = 0
    current_slide: int | None = None

    for b in blocks:
        if b.get("kind") not in ("text", "vision_placeholder"):
            continue
        content = str(b.get("content", ""))
        slide = b.get("slide")

        crosses_slide = slide is not None and current_slide is not None and slide != current_slide
        over_budget = current and (current_len + len(content) > char_budget)

        if current and (crosses_slide or over_budget):
            batches.append(current)
            current = []
            current_len = 0
            current_slide = None

        current.append(b)
        current_len += len(content)
        if slide is not None:
            current_slide = slide

    if current:
        batches.append(current)
    return batches


def _call_normalization_llm(combined: str, doc_type: str = "document") -> tuple[str, str | None]:
    """Returns (output_text, error_or_None). Azure → Gemini → Groq → passthrough."""
    if doc_type in ("slide_pdf", "pptx"):
        prompt = (
            "The following is text extracted from lecture slides. Reformat it as clean Markdown:\n"
            "- The first short line of each segment is the slide title → ## heading\n"
            "- Format bullet points with -\n"
            "- Remove standalone page numbers (a lone digit on its own line)\n"
            "- Remove navigation headers repeated across slides (e.g. repeated table-of-contents lines)\n"
            "- If formula characters are garbled or fragmented, replace with [Formula: brief description]\n"
            "- Preserve tables as Markdown tables\n"
            "Output clean Markdown. One ## section per slide.\n\n" + combined
        )
    else:
        prompt = (
            "Normalize the following extracted document text into clean, coherent Markdown:\n"
            "- Identify headings and use # / ## notation\n"
            "- Reconstruct table-like data as Markdown tables\n"
            "- Merge paragraph fragments broken across lines\n"
            "- Remove page headers/footers (repeated titles, standalone page numbers)\n"
            "- Preserve exact meaning, formulas, and structured data\n\n" + combined
        )

    if HAS_OPENAI:
        result = _call_azure_openai(prompt, operation="normalization")
        if not result.startswith("[LLM_ERROR]"):
            return result, None

    if HAS_GENAI and os.getenv("GEMINI_API_KEY"):
        try:
            return _call_gemini_vlm(prompt, operation="normalization_fallback_gemini"), None
        except Exception:
            pass

    groq_result = _call_groq(prompt, operation="normalization_fallback_groq")
    if groq_result.startswith("[LLM_ERROR]"):
        return combined, groq_result
    return groq_result, None


def llm_based_normalization(blocks: list[dict[str, Any]], doc_type: str = "document") -> tuple[str, list[dict[str, Any]]]:
    """Returns (normalized_markdown, per_batch_trace)."""
    batches = batch_blocks_for_normalization(blocks)
    if not batches:
        return "[EMPTY_OUTPUT]", []

    real_vision = os.getenv("OCR_ENABLE_REAL_VISION", "false").lower() == "true"

    outputs: list[str] = []
    trace: list[dict[str, Any]] = []

    for i, batch in enumerate(batches):
        combined = "\n\n".join(str(b.get("content", "")) for b in batch)
        page_range = [b.get("page") or b.get("slide") for b in batch if b.get("page") or b.get("slide")]
        trace_entry: dict[str, Any] = {
            "batch_index": i,
            "block_count": len(batch),
            "input_chars": len(combined),
            "page_or_slide_range": [page_range[0], page_range[-1]] if page_range else None,
        }

        if not combined.strip():
            trace_entry["mode"] = "skipped_empty"
            trace.append(trace_entry)
            continue

        if not real_vision:
            outputs.append(combined)
            trace_entry["mode"] = "passthrough_vision_disabled"
            trace_entry["output_chars"] = len(combined)
            trace.append(trace_entry)
            continue

        output, err = _call_normalization_llm(combined, doc_type)
        outputs.append(output)
        trace_entry["mode"] = "llm" if err is None else "llm_fallback_to_input"
        trace_entry["output_chars"] = len(output)
        if err:
            trace_entry["error"] = err
        trace.append(trace_entry)

    return "\n\n".join(outputs), trace


# ── Extraction ────────────────────────────────────────────────────────────────

def extract_pdf_text(
    pdf_path: Path,
) -> tuple[list[dict[str, Any]], list[dict[str, Any]], list[dict[str, Any]]]:
    """Returns (blocks, prompts, visual_trace)."""
    blocks: list[dict[str, Any]] = []
    prompts: list[dict[str, Any]] = []
    visual_trace: list[dict[str, Any]] = []

    try:
        import fitz  # PyMuPDF
    except Exception as exc:
        return blocks, [{"kind": "error", "message": f"PyMuPDF unavailable: {exc}"}], visual_trace

    enable_real_vision = os.getenv("OCR_ENABLE_REAL_VISION", "false").lower() == "true"
    inferred_type, _ = classify_pdf_type(pdf_path)

    doc = fitz.open(pdf_path)
    page_count = len(doc)

    # Pre-pass: images repeating across many pages are almost certainly decorative
    # (logos, banners, page borders). Hash-based detection avoids a VLM call entirely.
    image_hash_pages: dict[str, set[int]] = {}
    for page_index, page in enumerate(doc):
        for img_info in page.get_images(full=True):
            xref = img_info[0]
            try:
                base = doc.extract_image(xref)
                h = hashlib.sha256(base["image"]).hexdigest()
                image_hash_pages.setdefault(h, set()).add(page_index)
            except Exception:
                continue
    repeat_threshold = max(2, int(page_count * 0.3)) if page_count else 2
    decorative_repeat_hashes = {h for h, pages in image_hash_pages.items() if len(pages) >= repeat_threshold}

    for page_index, page in enumerate(doc):
        # Scanned PDF or slide PDF: render each page as an image for OCR/VLM.
        if inferred_type in ("scanned_pdf", "slide_pdf"):
            if enable_real_vision:
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                img_data = pix.tobytes("png")
                if inferred_type == "slide_pdf":
                    prompt_text = (
                        "This is a lecture slide page. Extract its content as Markdown:\n"
                        "- The slide title (large text at top) → ## heading\n"
                        "- Bullet points → - list items\n"
                        "- Any table → Markdown table\n"
                        "- Any formula or equation → LaTeX inside $...$\n"
                        "- Any diagram/graph/memory layout/flowchart → draw/represent it visually using a Markdown table (e.g. for variables/states/transitions) or a text-based ASCII diagram (using boxes [+---+] and arrows [-->] inside a preformatted ``` block). Underneath this visual representation, add a brief description in [Diagram: ...]\n"
                        "Output only the Markdown content, no commentary."
                    )
                else:
                    prompt_text = (
                        "This is a scanned mixed-layout page. Perform region-based OCR: "
                        "1. Detect table-like regions and extract as structured JSON. "
                        "2. Detect graph/diagram regions and extract as structured JSON. "
                        "3. Detect handwritten annotations and extract separately. "
                        "4. Extract body text. "
                        "Merge them back in approximate reading order."
                    )
                content = call_vlm(prompt_text, img_data)
                blocks.append({
                    "kind": "text",
                    "page": page_index + 1,
                    "content": content,
                    "ocr_strategy": "region_based",
                })
                prompts.append({
                    "kind": "vision_ocr",
                    "page": page_index + 1,
                    "instruction": prompt_text,
                    "ocr_strategy": "region_based",
                })
            else:
                blocks.append({
                    "kind": "vision_placeholder",
                    "page": page_index + 1,
                    "content": "[VISION_PLACEHOLDER] OCR output pending for scanned page.",
                })
                prompts.append({
                    "kind": "vision_prompt",
                    "page": page_index + 1,
                    "instruction": "Perform OCR on this scanned page.",
                })
            continue

        # Normal / mixed PDF: extract text layer first.
        text = (page.get_text("text") or "").strip()
        if text:
            blocks.append({"kind": "text", "page": page_index + 1, "content": text})

        page_height = page.rect.height if page.rect else None

        for image_idx, img_info in enumerate(page.get_images(full=True)):
            xref = img_info[0]
            if enable_real_vision:
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_hash = hashlib.sha256(image_bytes).hexdigest()

                try:
                    rects = page.get_image_rects(xref)
                    rect = rects[0] if rects else None
                except Exception:
                    rect = None
                position_hint = _image_position_hint(
                    rect.y0 if rect else None,
                    rect.y1 if rect else None,
                    page_height,
                )

                if image_hash in decorative_repeat_hashes:
                    classification = _make_classification(
                        "decorative", "none", 1.0, "repeated_across_pages", position_hint
                    )
                else:
                    classification = classify_visual(
                        image_bytes, surrounding_text=text, position_hint=position_hint
                    )

                trace_entry: dict[str, Any] = {
                    "source": "page",
                    "page": page_index + 1,
                    "image_index": image_idx,
                    "image_hash_prefix": image_hash[:16],
                    "repeat_count": len(image_hash_pages.get(image_hash, set())),
                    "position_hint": position_hint,
                    "classification": classification,
                    "action_taken": classification["action"],
                }
                visual_trace.append(trace_entry)

                if classification["action"] == "skip":
                    prompts.append({
                        "kind": "vision_skipped_decorative",
                        "page": page_index + 1,
                        "image_index": image_idx,
                        "position_hint": position_hint,
                        "repeat_count": trace_entry["repeat_count"],
                        "reason": classification["reason"],
                    })
                    continue

                if classification["action"] == "minimal_tag":
                    blocks.append({
                        "kind": "text",
                        "page": page_index + 1,
                        "image_index": image_idx,
                        "content": f"[Visual: decorative element, page {page_index + 1}]",
                    })
                    prompts.append({
                        "kind": "vision_minimal_tag",
                        "page": page_index + 1,
                        "image_index": image_idx,
                        "position_hint": position_hint,
                        "classification": classification,
                    })
                    continue

                # action == "extract"
                prompt_text = get_specialized_prompt(classification["label"])
                content = call_vlm(prompt_text, image_bytes)
                blocks.append({
                    "kind": "text",
                    "page": page_index + 1,
                    "image_index": image_idx,
                    "content": f"\n### Visual Element (Page {page_index + 1}, Image {image_idx})\n{content}\n",
                })
                prompts.append({
                    "kind": "vision_prompt",
                    "page": page_index + 1,
                    "image_index": image_idx,
                    "category": classification["label"],
                    "position_hint": position_hint,
                    "instruction": prompt_text,
                })
            else:
                prompts.append({
                    "kind": "vision_prompt",
                    "page": page_index + 1,
                    "image_index": image_idx,
                    "instruction": "Describe diagram/table/formula and convert to markdown text.",
                })
                blocks.append({
                    "kind": "vision_placeholder",
                    "page": page_index + 1,
                    "content": "[VISION_PLACEHOLDER] Describe visual element on this page.",
                })

    doc.close()
    return blocks, prompts, visual_trace


def extract_pptx_text(
    pptx_path: Path,
) -> tuple[list[dict[str, Any]], list[dict[str, Any]], list[dict[str, Any]]]:
    """Returns (blocks, prompts, visual_trace)."""
    blocks: list[dict[str, Any]] = []
    prompts: list[dict[str, Any]] = []
    visual_trace: list[dict[str, Any]] = []

    try:
        from pptx import Presentation
    except Exception as exc:
        return blocks, [{"kind": "error", "message": f"python-pptx unavailable: {exc}"}], visual_trace

    enable_real_vision = os.getenv("OCR_ENABLE_REAL_VISION", "false").lower() == "true"
    presentation = Presentation(str(pptx_path))

    # Pre-pass: detect images repeating across slides (logos, slide masters, banners).
    slide_total = len(presentation.slides)
    image_hash_slides: dict[str, set[int]] = {}
    for s_idx, slide in enumerate(presentation.slides):
        for shape in slide.shapes:
            is_picture = (
                getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE
                if HAS_PPTX_ENUMS
                else getattr(shape, "shape_type", None) == 13
            )
            if is_picture:
                try:
                    h = hashlib.sha256(shape.image.blob).hexdigest()
                    image_hash_slides.setdefault(h, set()).add(s_idx)
                except Exception:
                    continue
    repeat_threshold = max(2, int(slide_total * 0.3)) if slide_total else 2
    decorative_repeat_hashes = {h for h, slides in image_hash_slides.items() if len(slides) >= repeat_threshold}

    for slide_idx, slide in enumerate(presentation.slides):
        slide_text_parts: list[str] = []
        image_blocks: list[dict[str, Any]] = []

        # First pass: gather slide text so the vision classifier gets context.
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text_parts.append(shape.text.strip())
        slide_text_context = "\n".join(part for part in slide_text_parts if part)

        slide_height = presentation.slide_height or 0

        for shape in slide.shapes:
            is_picture = (
                getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE
                if HAS_PPTX_ENUMS
                else getattr(shape, "shape_type", None) == 13
            )

            if is_picture:
                if enable_real_vision:
                    image_bytes = shape.image.blob
                    image_hash = hashlib.sha256(image_bytes).hexdigest()

                    top = getattr(shape, "top", None)
                    height = getattr(shape, "height", None)
                    if top is not None and height is not None and slide_height:
                        position_hint = _image_position_hint(float(top), float(top + height), float(slide_height))
                    else:
                        position_hint = "inline"

                    if image_hash in decorative_repeat_hashes:
                        classification = _make_classification(
                            "decorative", "none", 1.0, "repeated_across_slides", position_hint
                        )
                    else:
                        classification = classify_visual(
                            image_bytes, surrounding_text=slide_text_context, position_hint=position_hint
                        )

                    trace_entry = {
                        "source": "slide",
                        "slide": slide_idx + 1,
                        "image_hash_prefix": image_hash[:16],
                        "repeat_count": len(image_hash_slides.get(image_hash, set())),
                        "position_hint": position_hint,
                        "classification": classification,
                        "action_taken": classification["action"],
                    }
                    visual_trace.append(trace_entry)

                    if classification["action"] == "skip":
                        prompts.append({
                            "kind": "vision_skipped_decorative",
                            "slide": slide_idx + 1,
                            "position_hint": position_hint,
                            "repeat_count": trace_entry["repeat_count"],
                            "reason": classification["reason"],
                        })
                        continue

                    if classification["action"] == "minimal_tag":
                        image_blocks.append({
                            "kind": "text",
                            "slide": slide_idx + 1,
                            "content": f"[Visual: decorative element, slide {slide_idx + 1}]",
                        })
                        prompts.append({
                            "kind": "vision_minimal_tag",
                            "slide": slide_idx + 1,
                            "position_hint": position_hint,
                            "classification": classification,
                        })
                        continue

                    # action == "extract"
                    prompt_text = get_specialized_prompt(classification["label"])
                    content = call_vlm(prompt_text, image_bytes)
                    image_blocks.append({
                        "kind": "text",
                        "slide": slide_idx + 1,
                        "content": f"\n### Slide Image\n{content}\n",
                    })
                    prompts.append({
                        "kind": "vision_prompt",
                        "slide": slide_idx + 1,
                        "category": classification["label"],
                        "position_hint": position_hint,
                        "instruction": prompt_text,
                    })
                else:
                    image_blocks.append({
                        "kind": "vision_placeholder",
                        "slide": slide_idx + 1,
                        "content": "[VISION_PLACEHOLDER] Describe embedded image.",
                    })
                    prompts.append({
                        "kind": "vision_prompt",
                        "slide": slide_idx + 1,
                        "instruction": "Describe embedded image and convert to markdown notes.",
                    })

        combined_text = "\n".join(part for part in slide_text_parts if part)
        if combined_text:
            blocks.append({"kind": "text", "slide": slide_idx + 1, "content": combined_text})

        blocks.extend(image_blocks)

    return blocks, prompts, visual_trace


def extract_image_text(
    image_path: Path,
) -> tuple[list[dict[str, Any]], list[dict[str, Any]], list[dict[str, Any]]]:
    """Returns (blocks, prompts, visual_trace)."""
    visual_trace: list[dict[str, Any]] = []
    enable_real_vision = os.getenv("OCR_ENABLE_REAL_VISION", "false").lower() == "true"

    if enable_real_vision:
        with open(image_path, "rb") as f:
            image_bytes = f.read()
        classification = classify_visual(image_bytes)
        visual_trace.append({
            "source": "image_file",
            "file": image_path.name,
            "classification": classification,
            "action_taken": classification["action"],
        })
        prompt_text = get_specialized_prompt(classification["label"]) + " Also perform OCR on any text."
        content = call_vlm(prompt_text, image_bytes)
        blocks = [{"kind": "text", "content": content}]
        prompts = [{"kind": "vision_ocr", "instruction": prompt_text, "image_path": image_path.name}]
    else:
        prompt_text = "Perform OCR and describe any table/formula/diagram in markdown format."
        prompts = [{"kind": "vision_prompt", "instruction": prompt_text, "image_path": image_path.name}]
        blocks = [{"kind": "vision_placeholder", "content": "[VISION_PLACEHOLDER] OCR output pending from vision model."}]

    return blocks, prompts, visual_trace


# ── Markdown assembly ─────────────────────────────────────────────────────────

def merge_blocks_to_markdown(sample_id: str, source_rel_path: str, blocks: list[dict[str, Any]]) -> str:
    lines = [f"# OCR Result: {sample_id}", "", f"Source: `{source_rel_path}`", ""]
    for block in blocks:
        if block["kind"] == "text":
            location = block.get("page") or block.get("slide") or "?"
            lines.append(f"## Segment {location}")
            lines.append("")
            lines.append(block.get("content", "").strip())
            lines.append("")
        elif block["kind"] == "vision_placeholder":
            location = block.get("page") or block.get("slide") or "?"
            lines.append(f"> Vision segment {location}: {block.get('content', '')}")
            lines.append("")

    content = "\n".join(lines).strip()
    if not content:
        return "# OCR Result\n\n[EMPTY_OUTPUT]"
    return content + "\n"


# ── Manifest validation ───────────────────────────────────────────────────────

@dataclass
class ValidationIssue:
    sample_id: str
    issue: str


def validate_manifest(project_root: Path, manifest_path: Path) -> tuple[list[dict[str, Any]], list[ValidationIssue]]:
    payload = read_json(manifest_path)
    entries = payload if isinstance(payload, list) else payload.get("entries", [])

    valid_entries: list[dict[str, Any]] = []
    issues: list[ValidationIssue] = []

    for idx, entry in enumerate(entries):
        sample_id = entry.get("id") or f"manifest_index_{idx}"
        required = ["id", "course_code", "tier", "input_path", "expected_type", "expected_route"]
        missing = [field for field in required if not entry.get(field)]
        if missing:
            issues.append(ValidationIssue(sample_id=sample_id, issue=f"missing_fields:{','.join(missing)}"))
            continue

        if entry["expected_type"] not in EXPECTED_TYPES:
            issues.append(ValidationIssue(sample_id=sample_id, issue="invalid_expected_type"))
            continue

        if entry["expected_route"] not in EXPECTED_ROUTES:
            issues.append(ValidationIssue(sample_id=sample_id, issue="invalid_expected_route"))
            continue

        input_path = project_root / entry["input_path"]
        if not input_path.exists() or not input_path.is_file():
            issues.append(ValidationIssue(sample_id=sample_id, issue="input_file_not_found"))
            continue

        valid_entries.append(entry)

    return valid_entries, issues


# ── Structured output rendering ───────────────────────────────────────────────

def render_deterministic_markdown(data: dict[str, Any]) -> str:
    content_type = data.get("content_type", "")
    if content_type in ("adjacency_matrix", "incidence_matrix", "table", "table_or_matrix"):
        rows = data.get("values", [])
        col_labels = data.get("column_labels", [])
        row_labels = data.get("row_labels", [])

        md_lines = []
        if col_labels:
            header = "| " + " | ".join([""] + [str(c) for c in col_labels]) + " |"
            sep = "| " + " | ".join(["---"] * (len(col_labels) + 1)) + " |"
            md_lines.extend([header, sep])
        elif rows and len(rows[0]) > 0:
            header = "| " + " | ".join(["Col " + str(i + 1) for i in range(len(rows[0]))]) + " |"
            sep = "| " + " | ".join(["---"] * len(rows[0])) + " |"
            md_lines.extend([header, sep])

        for i, row in enumerate(rows):
            r_label = str(row_labels[i]) if i < len(row_labels) else ""
            row_vals = [str(v) for v in row]
            if col_labels:
                line = "| " + " | ".join([r_label] + row_vals) + " |"
            else:
                line = "| " + " | ".join(row_vals) + " |"
            md_lines.append(line)

        if data.get("notes"):
            md_lines.append("")
            for note in data["notes"]:
                md_lines.append(f"**Note:** {note}")

        return "\n".join(md_lines)

    elif content_type == "graph_diagram":
        nodes = data.get("nodes", [])
        edges = data.get("edges", [])
        md_lines = ["**Graph/Diagram:**\n", "- Nodes: " + ", ".join(map(str, nodes))]
        if edges:
            md_lines.append("- Edges:")
            for e in edges:
                md_lines.append(f"  - {e}")
        return "\n".join(md_lines)

    return json.dumps(data, indent=2)


# ── Run mode ──────────────────────────────────────────────────────────────────

def run_mode(
    project_root: Path,
    output_root: Path,
    manifest_path: Path,
    run_id: str | None,
    progress_callback: Callable[[dict[str, Any]], None] | None = None,
) -> dict[str, Any]:
    global _llm_call_records, _progress_callback, _progress_records
    _llm_call_records = []  # reset per run so records don't bleed across calls
    _progress_records = []
    _progress_callback = progress_callback

    started_at = time.time()
    resolved_run_id = run_id or datetime.now().strftime("%Y%m%d_%H%M%S")
    run_root = output_root / resolved_run_id
    ensure_dir(run_root)
    _emit_progress("run_start", run_id=resolved_run_id, manifest_path=manifest_path.as_posix())

    valid_entries, issues = validate_manifest(project_root, manifest_path)
    manifest_validation = {
        "generated_at": utc_now_iso(),
        "manifest_path": manifest_path.relative_to(project_root).as_posix(),
        "valid_count": len(valid_entries),
        "invalid_count": len(issues),
        "issues": [{"sample_id": it.sample_id, "issue": it.issue} for it in issues],
    }
    write_json(run_root / "manifest_validation.json", manifest_validation)
    _emit_progress(
        "manifest_validated",
        run_id=resolved_run_id,
        valid_count=len(valid_entries),
        invalid_count=len(issues),
    )

    errors: list[dict[str, Any]] = []
    outcomes: list[dict[str, Any]] = []
    per_sample_metrics: list[dict[str, Any]] = []

    for entry in valid_entries:
        sample_id = entry["id"]
        sample_root = run_root / sample_id
        ensure_dir(sample_root)

        input_path = project_root / entry["input_path"]
        stage_started = time.time()
        _emit_progress(
            "sample_start",
            run_id=resolved_run_id,
            sample_id=sample_id,
            input_path=entry["input_path"],
            expected_type=entry["expected_type"],
            expected_route=entry["expected_route"],
        )

        try:
            route_started = time.time()
            _emit_progress("route_start", run_id=resolved_run_id, sample_id=sample_id)
            route, route_evidence = decide_route(input_path)
            _emit_progress(
                "route_end",
                run_id=resolved_run_id,
                sample_id=sample_id,
                route=route,
                route_evidence=route_evidence,
                duration_ms=int((time.time() - route_started) * 1000),
            )

            extraction_started = time.time()
            _emit_progress(
                "extraction_start",
                run_id=resolved_run_id,
                sample_id=sample_id,
                suffix=input_path.suffix.lower(),
            )
            if input_path.suffix.lower() == ".pdf":
                blocks, prompts, visual_trace = extract_pdf_text(input_path)
            elif input_path.suffix.lower() == ".pptx":
                blocks, prompts, visual_trace = extract_pptx_text(input_path)
            elif input_path.suffix.lower() in IMAGE_EXTENSIONS:
                blocks, prompts, visual_trace = extract_image_text(input_path)
            else:
                blocks, prompts, visual_trace = [], [{"kind": "error", "message": "unsupported_extension"}], []
            _emit_progress(
                "extraction_end",
                run_id=resolved_run_id,
                sample_id=sample_id,
                block_count=len(blocks),
                prompt_count=len(prompts),
                visual_count=len(visual_trace),
                duration_ms=int((time.time() - extraction_started) * 1000),
            )

            quality_flags = {
                "route_match": route == entry["expected_route"],
                "non_empty_output": False,
                "warnings": [],
                "validators": {
                    "parseable": None,
                    "matrix_shape": None,
                    "binary_values": None,
                    "table_consistency": None,
                },
            }

            structured_data = {}
            for b in blocks:
                if b["kind"] == "text" and "schema_version" in str(b.get("content", "")):
                    quality_flags["validators"]["parseable"] = False
                    try:
                        match = re.search(r"\{.*\}", str(b.get("content", "")), re.DOTALL)
                        if match:
                            parsed = json.loads(match.group(0))
                            structured_data = parsed
                            quality_flags["validators"]["parseable"] = True

                            c_type = parsed.get("content_type", "")
                            if c_type in ("adjacency_matrix", "incidence_matrix"):
                                vals = parsed.get("values", [])
                                if vals:
                                    num_cols = len(vals[0])
                                    quality_flags["validators"]["matrix_shape"] = (len(vals) == num_cols)
                                    quality_flags["validators"]["binary_values"] = all(
                                        str(v).strip() in ("0", "1") for row in vals for v in row
                                    )
                            elif c_type in ("table", "table_or_matrix"):
                                vals = parsed.get("values", [])
                                if vals:
                                    num_cols = len(vals[0])
                                    quality_flags["validators"]["table_consistency"] = all(
                                        len(row) == num_cols for row in vals
                                    )

                            b["content"] = render_deterministic_markdown(parsed)
                    except Exception:
                        pass

            result_md = merge_blocks_to_markdown(sample_id, entry["input_path"], blocks)
            quality_flags["non_empty_output"] = bool(result_md.strip())

            write_text(sample_root / "raw_result.md", result_md)

            cleaned_blocks = rule_based_cleanup(blocks)
            ext = input_path.suffix.lower()
            if ext == ".pptx":
                doc_type = "pptx"
            elif ext == ".pdf":
                doc_type = route_evidence.get("inferred_type", "document")
            else:
                doc_type = "document"
            normalization_started = time.time()
            _emit_progress(
                "normalization_start",
                run_id=resolved_run_id,
                sample_id=sample_id,
                doc_type=doc_type,
                block_count=len(cleaned_blocks),
            )
            normalized_md, normalization_batches = llm_based_normalization(cleaned_blocks, doc_type)
            _emit_progress(
                "normalization_end",
                run_id=resolved_run_id,
                sample_id=sample_id,
                batch_count=len(normalization_batches),
                normalized_chars=len(normalized_md),
                duration_ms=int((time.time() - normalization_started) * 1000),
            )
            write_text(sample_root / "normalized_result.md", normalized_md)
            write_text(sample_root / "result.md", normalized_md)

            write_json(sample_root / "structured.json", structured_data)
            write_json(
                sample_root / "normalization_trace.json",
                {
                    "char_budget": NORMALIZATION_CHAR_BUDGET,
                    "batches": normalization_batches,
                    "cleaned_blocks": cleaned_blocks,
                },
            )

            # Visual classification trace per sample.
            vis_skipped = sum(1 for t in visual_trace if t["action_taken"] == "skip")
            vis_minimal = sum(1 for t in visual_trace if t["action_taken"] == "minimal_tag")
            vis_extracted = sum(1 for t in visual_trace if t["action_taken"] == "extract")
            write_json(
                sample_root / "visual_classification_trace.json",
                {
                    "sample_id": sample_id,
                    "total_images": len(visual_trace),
                    "decorative_skipped": vis_skipped,
                    "decorative_minimal_tagged": vis_minimal,
                    "meaningful_extracted": vis_extracted,
                    "entries": visual_trace,
                },
            )

            metadata = {
                "sample_id": sample_id,
                "course_code": entry["course_code"],
                "tier": entry["tier"],
                "input_path": entry["input_path"],
                "expected_type": entry["expected_type"],
                "expected_route": entry["expected_route"],
                "actual_route": route,
                "route_evidence": route_evidence,
                "file_sha256": file_sha256(input_path),
                "model": "local_placeholder",
                "duration_ms": int((time.time() - stage_started) * 1000),
            }
            write_json(sample_root / "metadata.json", metadata)
            write_json(sample_root / "vision_prompts.json", prompts)

            if not quality_flags["route_match"]:
                quality_flags["warnings"].append("route_mismatch")
            if "[VISION_PLACEHOLDER]" in result_md:
                quality_flags["warnings"].append("vision_output_placeholder_detected")

            is_valid = True
            if quality_flags["validators"].get("parseable") is False:
                is_valid = False
                quality_flags["warnings"].append("structured_content_parse_failed")

            for _, v in quality_flags["validators"].items():
                if v is False:
                    is_valid = False
                    if "validator_failed" not in quality_flags["warnings"]:
                        quality_flags["warnings"].append("validator_failed")

            write_json(sample_root / "quality_flags.json", quality_flags)

            # Per-sample metrics for run-level aggregation.
            raw_chars = sum(len(str(b.get("content", ""))) for b in blocks)
            normalized_chars = len(normalized_md)
            placeholder_count = result_md.count("[VISION_PLACEHOLDER]")
            validator_failures = sum(1 for v in quality_flags["validators"].values() if v is False)
            per_sample_metrics.append({
                "sample_id": sample_id,
                "total_images": len(visual_trace),
                "decorative_skipped": vis_skipped,
                "decorative_minimal_tagged": vis_minimal,
                "meaningful_extracted": vis_extracted,
                "raw_chars": raw_chars,
                "normalized_chars": normalized_chars,
                "placeholder_count": placeholder_count,
                "validator_failures": validator_failures,
            })

            outcomes.append({
                "sample_id": sample_id,
                "status": "success" if is_valid else "failed",
                "route_match": quality_flags["route_match"],
                "output_non_empty": quality_flags["non_empty_output"],
            })
            _emit_progress(
                "sample_end",
                run_id=resolved_run_id,
                sample_id=sample_id,
                status="success" if is_valid else "failed",
                duration_ms=int((time.time() - stage_started) * 1000),
            )

        except Exception as exc:
            _emit_progress(
                "sample_error",
                run_id=resolved_run_id,
                sample_id=sample_id,
                error_type=exc.__class__.__name__,
                message=str(exc),
                duration_ms=int((time.time() - stage_started) * 1000),
            )
            errors.append({
                "sample_id": sample_id,
                "stage": "run",
                "error_type": exc.__class__.__name__,
                "message": str(exc),
                "retry_count": 0,
                "timestamp": utc_now_iso(),
            })
            outcomes.append({
                "sample_id": sample_id,
                "status": "failed",
                "route_match": False,
                "output_non_empty": False,
            })
            per_sample_metrics.append({
                "sample_id": sample_id,
                "error": str(exc),
            })

    expected_coverage: dict[str, int] = {kind: 0 for kind in REQUIRED_EDGE_CASES}
    for entry in valid_entries:
        expected_coverage[entry["expected_type"]] = expected_coverage.get(entry["expected_type"], 0) + 1

    coverage_report = {
        "generated_at": utc_now_iso(),
        "expected_type_counts": expected_coverage,
        "covered_categories": [key for key, count in expected_coverage.items() if count > 0],
        "missing_categories": [key for key, count in expected_coverage.items() if count == 0],
    }

    success_count = sum(1 for outcome in outcomes if outcome["status"] == "success")
    failed_count = sum(1 for outcome in outcomes if outcome["status"] == "failed")
    route_match_count = sum(1 for outcome in outcomes if outcome.get("route_match"))

    if valid_entries and failed_count == 0:
        run_status = "success"
    elif success_count > 0:
        run_status = "partial_success"
    else:
        run_status = "failed"

    run_summary = {
        "generated_at": utc_now_iso(),
        "run_id": resolved_run_id,
        "status": run_status,
        "manifest_path": manifest_path.relative_to(project_root).as_posix(),
        "total_valid_entries": len(valid_entries),
        "total_invalid_entries": len(issues),
        "success_count": success_count,
        "failed_count": failed_count,
        "route_match_count": route_match_count,
        "duration_ms": int((time.time() - started_at) * 1000),
    }

    # Aggregate experiment metrics across all samples.
    def _sum_field(field: str) -> int:
        return sum(s.get(field, 0) for s in per_sample_metrics if isinstance(s.get(field), int))

    total_images = _sum_field("total_images")
    meaningful = _sum_field("meaningful_extracted")
    structured_attempts = sum(
        1 for s in per_sample_metrics
        if isinstance(s.get("validator_failures"), int)
    )
    structured_pass = sum(
        1 for s in per_sample_metrics
        if isinstance(s.get("validator_failures"), int) and s["validator_failures"] == 0
    )
    metrics = {
        "generated_at": utc_now_iso(),
        "run_id": resolved_run_id,
        "total_images_classified": total_images,
        "decorative_images_skipped": _sum_field("decorative_skipped"),
        "decorative_minimal_tagged": _sum_field("decorative_minimal_tagged"),
        "meaningful_visuals_extracted": meaningful,
        "decorative_ratio": round(
            (_sum_field("decorative_skipped") + _sum_field("decorative_minimal_tagged")) / max(total_images, 1), 3
        ),
        "raw_chars_total": _sum_field("raw_chars"),
        "normalized_chars_total": _sum_field("normalized_chars"),
        "placeholder_count_total": _sum_field("placeholder_count"),
        "validator_failure_total": _sum_field("validator_failures"),
        "structured_json_parse_pass_rate": round(structured_pass / max(structured_attempts, 1), 3),
        "per_sample": per_sample_metrics,
    }

    # LLM API cost/latency metrics for this run.
    llm_summary = _summarize_llm_metrics()
    write_json(run_root / "llm_metrics.json", llm_summary)
    # Mirror aggregate (no per-call records) into metrics.json as llm_api.
    metrics["llm_api"] = {k: v for k, v in llm_summary.items() if k != "records"}
    _emit_progress(
        "run_end",
        run_id=resolved_run_id,
        status=run_status,
        duration_ms=run_summary["duration_ms"],
        success_count=success_count,
        failed_count=failed_count,
    )

    write_json(run_root / "run_summary.json", run_summary)
    write_json(run_root / "coverage_report.json", coverage_report)
    write_json(run_root / "errors.json", errors)
    write_json(run_root / "metrics.json", metrics)
    write_json(run_root / "progress_trace.json", _progress_records)

    progress_trace = list(_progress_records)
    _progress_callback = None
    return {
        "run_root": run_root,
        "manifest_validation": manifest_validation,
        "run_summary": run_summary,
        "coverage_report": coverage_report,
        "errors": errors,
        "metrics": metrics,
        "llm_metrics": llm_summary,
        "progress_trace": progress_trace,
    }


# ── CLI ───────────────────────────────────────────────────────────────────────

def build_arg_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Document processing experimentation pipeline runner")
    parser.add_argument("mode", choices=["inventory", "gap-report", "run", "test-api"])
    parser.add_argument(
        "--project-root",
        default=None,
        help="Absolute or relative project root. Defaults to repository root.",
    )
    parser.add_argument(
        "--manifest",
        default="src/experiments/document_processing/dataset_manifest.json",
        help="Manifest path for run mode (relative to project root by default).",
    )
    parser.add_argument("--run-id", default=None, help="Optional run id for run mode output folder.")
    return parser


def test_api_mode() -> dict[str, Any]:
    """Connectivity check for all configured LLM providers. Returns per-provider status."""
    results: dict[str, Any] = {}

    # ── Azure ─────────────────────────────────────────────────────────────────
    endpoint = os.getenv("AZURE_OPENAI_ENDPOINT", _AZURE_ENDPOINT_DEFAULT)
    deployment = os.getenv("AZURE_OPENAI_DEPLOYMENT", _AZURE_DEPLOYMENT_DEFAULT)
    azure_key = os.getenv("AZURE_AI_API_KEY")
    if azure_key and HAS_OPENAI:
        t0 = time.time()
        try:
            client = OpenAI(base_url=endpoint, api_key=azure_key, timeout=15.0)
            resp = client.chat.completions.create(
                model=deployment,
                messages=[{"role": "user", "content": "Reply with the single word: ok"}],
                max_tokens=5,
                temperature=0,
            )
            latency_ms = int((time.time() - t0) * 1000)
            results["azure_openai"] = {
                "status": "ok",
                "latency_ms": latency_ms,
                "model": deployment,
                "reply": (resp.choices[0].message.content or "").strip(),
            }
        except Exception as e:
            latency_ms = int((time.time() - t0) * 1000)
            status_code = getattr(e, "status_code", None)
            results["azure_openai"] = {
                "status": "error",
                "latency_ms": latency_ms,
                "error_type": type(e).__name__,
                "http_status": status_code,
                "detail": str(e)[:400],
            }
    else:
        results["azure_openai"] = {
            "status": "skipped",
            "reason": "AZURE_AI_API_KEY missing" if not azure_key else "openai package not installed",
        }

    # ── Gemini ────────────────────────────────────────────────────────────────
    gemini_key = os.getenv("GEMINI_API_KEY")
    if gemini_key and HAS_GENAI:
        model_id = os.getenv("VLM_MODEL_ID", "gemini-3.1-flash-lite")
        t0 = time.time()
        try:
            client = genai.Client(api_key=gemini_key, http_options={"timeout": 15})
            resp = client.models.generate_content(
                model=model_id,
                contents=["Reply with the single word: ok"],
                config=types.GenerateContentConfig(max_output_tokens=5, temperature=0),
            )
            latency_ms = int((time.time() - t0) * 1000)
            results["gemini"] = {
                "status": "ok",
                "latency_ms": latency_ms,
                "model": model_id,
                "reply": (resp.text or "").strip(),
            }
        except Exception as e:
            latency_ms = int((time.time() - t0) * 1000)
            status_code = getattr(e, "status_code", None) or getattr(e, "code", None)
            results["gemini"] = {
                "status": "error",
                "latency_ms": latency_ms,
                "error_type": type(e).__name__,
                "http_status": status_code,
                "detail": str(e)[:400],
            }
    else:
        results["gemini"] = {
            "status": "skipped",
            "reason": "GEMINI_API_KEY missing" if not gemini_key else "google-genai package not installed",
        }

    # ── Groq ──────────────────────────────────────────────────────────────────
    groq_key = os.getenv("GROQ_API_KEY")
    if groq_key and HAS_OPENAI:
        model = os.getenv("GROQ_MODEL_ID", _GROQ_MODEL_DEFAULT)
        t0 = time.time()
        try:
            client = OpenAI(api_key=groq_key, base_url=_GROQ_BASE_URL, timeout=15.0)
            resp = client.chat.completions.create(
                model=model,
                messages=[{"role": "user", "content": "Reply with the single word: ok"}],
                max_tokens=5,
                temperature=0,
            )
            latency_ms = int((time.time() - t0) * 1000)
            results["groq"] = {
                "status": "ok",
                "latency_ms": latency_ms,
                "model": model,
                "reply": (resp.choices[0].message.content or "").strip(),
            }
        except Exception as e:
            latency_ms = int((time.time() - t0) * 1000)
            status_code = getattr(e, "status_code", None)
            results["groq"] = {
                "status": "error",
                "latency_ms": latency_ms,
                "error_type": type(e).__name__,
                "http_status": status_code,
                "detail": str(e)[:400],
            }
    else:
        results["groq"] = {
            "status": "skipped",
            "reason": "GROQ_API_KEY missing" if not groq_key else "openai package not installed",
        }

    return results


def check_api_key() -> None:
    enable_real_vision = os.getenv("OCR_ENABLE_REAL_VISION", "false").lower() == "true"
    if not enable_real_vision:
        return
    has_azure = bool(os.getenv("AZURE_AI_API_KEY"))
    has_gemini = bool(os.getenv("GEMINI_API_KEY"))
    if not has_azure:
        print("WARNING: AZURE_AI_API_KEY missing — Azure is primary VLM, will fall back to Gemini.")
    if not has_gemini:
        print("WARNING: GEMINI_API_KEY missing — Gemini fallback unavailable.")
    if not has_azure and not has_gemini:
        print("WARNING: No VLM provider with image support. Groq (text-only) is last resort — vision quality degraded.")


def main() -> int:
    load_dotenv()
    load_dotenv(Path(__file__).parent.parent.parent / ".env", override=False)
    check_api_key()
    args = build_arg_parser().parse_args()
    project_root = resolve_project_root(args.project_root)
    output_root = resolve_output_root(project_root)
    ensure_dir(output_root)

    if args.mode == "test-api":
        results = test_api_mode()
        print(json.dumps({"mode": "test-api", "providers": results}, ensure_ascii=False, indent=2))
        return 0

    if args.mode == "inventory":
        report = inventory_mode(project_root, output_root)
        print(json.dumps({"mode": "inventory", "total_files": report["total_files"]}, ensure_ascii=False))
        return 0

    if args.mode == "gap-report":
        report = gap_report_mode(project_root, output_root)
        print(json.dumps({"mode": "gap-report", "missing_categories": report["missing_categories"]}, ensure_ascii=False))
        return 0

    manifest_path = Path(args.manifest)
    if not manifest_path.is_absolute():
        manifest_path = project_root / manifest_path
    if not manifest_path.exists():
        print(json.dumps({"error": "manifest_not_found", "manifest": str(manifest_path)}))
        return 2

    result = run_mode(project_root, output_root, manifest_path, args.run_id)
    print(
        json.dumps(
            {
                "mode": "run",
                "run_id": result["run_summary"]["run_id"],
                "status": result["run_summary"]["status"],
                "success_count": result["run_summary"]["success_count"],
                "failed_count": result["run_summary"]["failed_count"],
                "decorative_skipped": result["metrics"]["decorative_images_skipped"],
                "meaningful_extracted": result["metrics"]["meaningful_visuals_extracted"],
            },
            ensure_ascii=False,
        )
    )
    return 0


if __name__ == "__main__":
    sys.exit(main())
