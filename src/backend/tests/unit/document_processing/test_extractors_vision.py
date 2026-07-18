from __future__ import annotations

import io
import json
import threading
import time
from pathlib import Path

import fitz  # type: ignore[import-untyped]
from PIL import Image
from pptx import Presentation

from app.services.document_processing.classification import VisualClassifier
from app.services.document_processing.config import DocumentProcessingConfig
from app.services.document_processing.extractors.pdf import PdfExtractor
from app.services.document_processing.extractors.pptx import PptxExtractor
from app.services.document_processing.providers.base import (
    ProviderResponse,
    VisionLanguageProvider,
)

from .conftest import make_chain


def _png_bytes(color: tuple[int, int, int]) -> bytes:
    buf = io.BytesIO()
    Image.new("RGB", (40, 40), color).save(buf, format="PNG")
    return buf.getvalue()


class _ConcurrencyTrackingProvider(VisionLanguageProvider):
    """Stub VLM: classify() calls get JSON, other calls get plain text.

    Sleeps briefly per call and records concurrent-in-flight count so tests can assert
    that calls actually overlap instead of running sequentially.
    """

    provider_name = "stub"

    def __init__(self, *, sleep_seconds: float = 0.05) -> None:
        self._sleep_seconds = sleep_seconds
        self.call_count = 0
        self.max_concurrent = 0
        self._in_flight = 0
        self._lock = threading.Lock()

    def complete(
        self,
        prompt: str,
        *,
        images: list[bytes] | None = None,
        operation: str = "text",
    ) -> ProviderResponse:
        with self._lock:
            self.call_count += 1
            self._in_flight += 1
            self.max_concurrent = max(self.max_concurrent, self._in_flight)
        try:
            time.sleep(self._sleep_seconds)
            if "Classify this image" in prompt:
                text = json.dumps(
                    {"label": "chart_plot", "learning_value": "high", "confidence": 0.9, "reason": "test"}
                )
            else:
                text = f"extracted:{prompt[:20]}"
            return ProviderResponse(
                text=text, provider="stub", model="stub", status="success", latency_ms=0
            )
        finally:
            with self._lock:
                self._in_flight -= 1


def test_pptx_vision_extraction_runs_concurrently_and_preserves_order(tmp_path: Path) -> None:
    pptx_path = tmp_path / "deck.pptx"
    presentation = Presentation()
    image_paths = []
    for i in range(6):
        img_path = tmp_path / f"img{i}.png"
        img_path.write_bytes(_png_bytes((i * 10, 0, 0)))
        image_paths.append(img_path)

    for img_path in image_paths:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_picture(str(img_path), 0, 0, width=100, height=100)
    presentation.save(str(pptx_path))

    provider = _ConcurrencyTrackingProvider(sleep_seconds=0.05)
    config = DocumentProcessingConfig(enable_real_vision=True, vision_max_workers=4)
    chain = make_chain([provider], enable_real_vision=True)
    classifier = VisualClassifier(config, provider_chain=chain)

    from app.services.document_processing.progress import ProgressEmitter

    extractor = PptxExtractor(config, provider_chain=chain, classifier=classifier, emitter=ProgressEmitter())
    result = extractor.extract(pptx_path)

    # Each of the 6 unique images needs a classify call + an extract call (chart_plot -> extract).
    assert provider.call_count == 12
    # With max_workers=4 and unique per-image hashes, calls should overlap.
    assert provider.max_concurrent > 1

    # Slide image blocks appear in slide order.
    slide_numbers = [
        b["slide"]
        for b in result.blocks
        if b.get("kind") == "text" and "Slide Image" in b.get("content", "")
    ]
    assert slide_numbers == sorted(slide_numbers)
    assert len(slide_numbers) == 6


def test_pdf_embedded_image_extraction_preserves_order(tmp_path: Path) -> None:
    pdf_path = tmp_path / "doc.pdf"
    # Enough body text per page (avg_chars_per_page >= 80, non-landscape) so RouteDecider
    # classifies this as "mixed_pdf" and takes the per-embedded-image extraction path
    # rather than the whole-page-render path used for scanned/slide PDFs.
    body_text = "Lorem ipsum dolor sit amet consectetur adipiscing elit. " * 3
    doc = fitz.open()
    for i in range(4):
        page = doc.new_page()
        page.insert_text((72, 72), f"Page {i}: {body_text}")
        page.insert_image(fitz.Rect(72, 300, 172, 400), stream=_png_bytes((0, i * 10, 0)))
    doc.save(str(pdf_path))
    doc.close()

    from app.services.document_processing.routing import RouteDecider

    inferred_type, _ = RouteDecider().classify_pdf_type(pdf_path)
    assert inferred_type not in ("scanned_pdf", "slide_pdf")

    provider = _ConcurrencyTrackingProvider(sleep_seconds=0.05)
    config = DocumentProcessingConfig(enable_real_vision=True, vision_max_workers=4)
    chain = make_chain([provider], enable_real_vision=True)
    classifier = VisualClassifier(config, provider_chain=chain)

    from app.services.document_processing.progress import ProgressEmitter

    extractor = PdfExtractor(config, provider_chain=chain, classifier=classifier, emitter=ProgressEmitter())
    result = extractor.extract(pdf_path)

    assert provider.call_count == 8  # 4 images x (classify + extract)
    assert provider.max_concurrent > 1

    page_numbers = [
        b["page"]
        for b in result.blocks
        if b.get("kind") == "text" and "Visual Element" in b.get("content", "")
    ]
    assert page_numbers == sorted(page_numbers)
    assert len(page_numbers) == 4


def test_pdf_scanned_page_render_extraction_runs_concurrently(tmp_path: Path) -> None:
    pdf_path = tmp_path / "scanned.pdf"
    doc = fitz.open()
    for i in range(4):
        page = doc.new_page()
        page.insert_image(fitz.Rect(0, 0, 200, 200), stream=_png_bytes((10 * i, 0, 0)))
    doc.save(str(pdf_path))
    doc.close()

    from app.services.document_processing.routing import RouteDecider

    inferred_type, _ = RouteDecider().classify_pdf_type(pdf_path)
    assert inferred_type == "scanned_pdf"

    provider = _ConcurrencyTrackingProvider(sleep_seconds=0.05)
    config = DocumentProcessingConfig(enable_real_vision=True, vision_max_workers=4)
    chain = make_chain([provider], enable_real_vision=True)
    classifier = VisualClassifier(config, provider_chain=chain)

    from app.services.document_processing.progress import ProgressEmitter

    extractor = PdfExtractor(config, provider_chain=chain, classifier=classifier, emitter=ProgressEmitter())
    result = extractor.extract(pdf_path)

    # One whole-page vision call per page (no classify step for rendered pages).
    assert provider.call_count == 4
    assert provider.max_concurrent > 1

    page_numbers = [b["page"] for b in result.blocks if b.get("kind") == "text"]
    assert page_numbers == [1, 2, 3, 4]
