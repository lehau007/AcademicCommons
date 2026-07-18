from __future__ import annotations

from pathlib import Path

import fitz  # type: ignore[import-untyped]
from pptx import Presentation

from app.services.document_processing import (
    DocumentProcessingConfig,
    DocumentProcessingPipeline,
)


def test_pipeline_text_pdf_end_to_end_offline(tmp_path: Path) -> None:
    pdf_path = tmp_path / "doc.pdf"
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Hello " * 200)
    doc.save(str(pdf_path))
    doc.close()

    pipeline = DocumentProcessingPipeline(DocumentProcessingConfig())
    result = pipeline.process_document(
        pdf_path, document_id="doc-1", expected_route="direct_text"
    )

    assert result.route == "direct_text"
    assert "Hello" in result.markdown
    assert result.quality_flags["route_match"] is True
    # No LLM calls happen when vision is disabled.
    assert result.llm_metrics["total_calls"] == 0


def test_pipeline_pptx_end_to_end_offline(tmp_path: Path) -> None:
    pptx_path = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Slide Title"
    slide.placeholders[1].text = "Slide body content for testing."
    presentation.save(str(pptx_path))

    pipeline = DocumentProcessingPipeline(DocumentProcessingConfig())
    result = pipeline.process_document(
        pptx_path, document_id="deck-1", expected_route="hybrid"
    )

    assert result.route == "hybrid"
    assert result.markdown.strip()
    assert result.llm_metrics["total_calls"] == 0
