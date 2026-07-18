"""PPTX extractor. Ports experiment ``extract_pptx_text`` (lines ~1132-1278)."""

from __future__ import annotations

import hashlib
from concurrent.futures import ThreadPoolExecutor
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from app.services.document_processing.classification import VisualClassifier
from app.services.document_processing.extractors.base import Extractor
from app.services.document_processing.models import ExtractionResult


@dataclass
class _PictureTask:
    slide_idx: int
    image_bytes: bytes
    image_hash: str
    position_hint: str
    slide_text_context: str
    repeat_count: int
    is_decorative_repeat: bool


class PptxExtractor(Extractor):
    def extract(self, path: Path) -> ExtractionResult:
        blocks: list[dict[str, Any]] = []
        prompts: list[dict[str, Any]] = []
        visual_trace: list[dict[str, Any]] = []

        try:
            from pptx import Presentation
        except Exception as exc:
            return ExtractionResult(
                blocks=blocks,
                prompts=[{"kind": "error", "message": f"python-pptx unavailable: {exc}"}],
                visual_trace=visual_trace,
            )

        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            has_pptx_enums = True
        except ImportError:
            has_pptx_enums = False

        def _is_picture(shape: Any) -> bool:
            if has_pptx_enums:
                return getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE
            return getattr(shape, "shape_type", None) == 13

        enable_real_vision = self._config.enable_real_vision
        presentation = Presentation(str(path))

        # Pre-pass: detect images repeating across slides (logos, slide masters, banners).
        slide_total = len(presentation.slides)
        image_hash_slides: dict[str, set[int]] = {}
        for s_idx, slide in enumerate(presentation.slides):
            for shape in slide.shapes:
                if _is_picture(shape):
                    try:
                        h = hashlib.sha256(shape.image.blob).hexdigest()
                        image_hash_slides.setdefault(h, set()).add(s_idx)
                    except Exception:
                        continue
        repeat_threshold = max(2, int(slide_total * 0.3)) if slide_total else 2
        decorative_repeat_hashes = {
            h for h, slides in image_hash_slides.items() if len(slides) >= repeat_threshold
        }

        # Pass 1 (no network calls): gather slide text and enumerate picture shapes.
        slide_combined_text: dict[int, str] = {}
        image_blocks_by_slide: dict[int, list[dict[str, Any]]] = {}
        picture_tasks: list[_PictureTask] = []
        slide_height = presentation.slide_height or 0

        for slide_idx, slide in enumerate(presentation.slides):
            slide_text_parts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text_parts.append(shape.text.strip())
            slide_text_context = "\n".join(part for part in slide_text_parts if part)
            slide_combined_text[slide_idx] = slide_text_context

            for shape in slide.shapes:
                if not _is_picture(shape):
                    continue

                if not enable_real_vision:
                    image_blocks_by_slide.setdefault(slide_idx, []).append({
                        "kind": "vision_placeholder",
                        "slide": slide_idx + 1,
                        "content": "[VISION_PLACEHOLDER] Describe embedded image.",
                    })
                    prompts.append({
                        "kind": "vision_prompt",
                        "slide": slide_idx + 1,
                        "instruction": "Describe embedded image and convert to markdown notes.",
                    })
                    continue

                image_bytes = shape.image.blob
                image_hash = hashlib.sha256(image_bytes).hexdigest()
                top = getattr(shape, "top", None)
                height = getattr(shape, "height", None)
                if top is not None and height is not None and slide_height:
                    position_hint = self._classifier.position_hint(
                        float(top), float(top + height), float(slide_height)
                    )
                else:
                    position_hint = "inline"

                picture_tasks.append(
                    _PictureTask(
                        slide_idx=slide_idx,
                        image_bytes=image_bytes,
                        image_hash=image_hash,
                        position_hint=position_hint,
                        slide_text_context=slide_text_context,
                        repeat_count=len(image_hash_slides.get(image_hash, set())),
                        is_decorative_repeat=image_hash in decorative_repeat_hashes,
                    )
                )

        # Pass 2 (network calls, parallelized): classify + extract each picture. Vision
        # calls dominate runtime on image-heavy decks, so this is the main OCR bottleneck.
        def _run_task(task: _PictureTask) -> dict[str, Any]:
            slide_idx = task.slide_idx
            if task.is_decorative_repeat:
                classification = self._classifier.make_classification(
                    "decorative", "none", 1.0, "repeated_across_slides", task.position_hint
                )
            else:
                classification = self._classifier.classify(
                    task.image_bytes,
                    surrounding_text=task.slide_text_context,
                    position_hint=task.position_hint,
                )

            trace_entry: dict[str, Any] = {
                "source": "slide",
                "slide": slide_idx + 1,
                "image_hash_prefix": task.image_hash[:16],
                "repeat_count": task.repeat_count,
                "position_hint": task.position_hint,
                "classification": classification,
                "action_taken": classification["action"],
            }

            if classification["action"] == "skip":
                return {
                    "trace_entry": trace_entry,
                    "prompt": {
                        "kind": "vision_skipped_decorative",
                        "slide": slide_idx + 1,
                        "position_hint": task.position_hint,
                        "repeat_count": task.repeat_count,
                        "reason": classification["reason"],
                    },
                    "block": None,
                }

            if classification["action"] == "minimal_tag":
                return {
                    "trace_entry": trace_entry,
                    "prompt": {
                        "kind": "vision_minimal_tag",
                        "slide": slide_idx + 1,
                        "position_hint": task.position_hint,
                        "classification": classification,
                    },
                    "block": {
                        "kind": "text",
                        "slide": slide_idx + 1,
                        "content": f"[Visual: decorative element, slide {slide_idx + 1}]",
                    },
                }

            # action == "extract"
            prompt_text = VisualClassifier.specialized_prompt(
                classification["label"], context=task.slide_text_context
            )
            content = self._chain.vision(prompt_text, task.image_bytes)
            return {
                "trace_entry": trace_entry,
                "prompt": {
                    "kind": "vision_prompt",
                    "slide": slide_idx + 1,
                    "category": classification["label"],
                    "position_hint": task.position_hint,
                    "instruction": prompt_text,
                },
                "block": {
                    "kind": "text",
                    "slide": slide_idx + 1,
                    "content": f"\n### Slide Image\n{content}\n",
                },
            }

        if picture_tasks:
            max_workers = max(1, self._config.vision_max_workers)
            with ThreadPoolExecutor(max_workers=max_workers) as executor:
                for task, result in zip(picture_tasks, executor.map(_run_task, picture_tasks), strict=True):
                    visual_trace.append(result["trace_entry"])
                    prompts.append(result["prompt"])
                    if result["block"] is not None:
                        image_blocks_by_slide.setdefault(task.slide_idx, []).append(result["block"])

        for slide_idx in range(slide_total):
            combined_text = slide_combined_text.get(slide_idx, "")
            if combined_text:
                blocks.append({"kind": "text", "slide": slide_idx + 1, "content": combined_text})
            blocks.extend(image_blocks_by_slide.get(slide_idx, []))

        return ExtractionResult(blocks=blocks, prompts=prompts, visual_trace=visual_trace)
